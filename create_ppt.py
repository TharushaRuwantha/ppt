from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

RED   = RGBColor(0xFF, 0x00, 0x00)
BLUE  = RGBColor(0x00, 0x00, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

questions_data = [
    {
        "number": "01",
        "question": (
            "A stone falls freely from rest and the total distance covered by it "
            "in the last second of its motion equals the distance covered by it in "
            "the first three seconds of its motion. The stone remains in the air for,"
        ),
        "options": ["(1)   6 s", "(2)   5 s", "(3)   7 s", "(4)   4 s", "(5)   3 s"],
        "answer": "Ans: (2)",
    },
    {
        "number": "02",
        "question": (
            "A car starting from rest and travelling with uniform acceleration obtains "
            "a velocity 5 m s\u207b\u00b9 in 2 s. After how many seconds will the car "
            "attain a velocity of 15 m s\u207b\u00b9?"
        ),
        "options": ["(1)   6 s", "(2)   4 s", "(3)   8 s", "(4)   15 s", "(5)   10 s"],
        "answer": "Ans: (2)",
    },
    {
        "number": "03",
        "question": (
            "A v-t graph of a moving particle is given below. What are the true "
            "statements about it?\n"
            "(A)  The particle is moving in the same direction at A and D\n"
            "(B)  The particle is decelerating at B and accelerating at F\n"
            "(C)  It changes the direction at C and E"
        ),
        "options": [
            "(1)   Only (A) and (B) are correct",
            "(2)   Only (A) and (C) are correct",
            "(3)   Only (B) and (C) are correct",
            "(4)   Only (C) and (D) are correct",
            "(5)   All (A), (B) and (C) are correct",
        ],
        "answer": "Ans: (3)",
    },
    {
        "number": "04",
        "question": (
            "The following statements are about a town X:\n"
            "(A)  The town is a five mile drive along the winding country road.\n"
            "(B)  The town sits at an altitude of 940 m.\n"
            "(C)  The town X is ten miles north of the town Y."
        ),
        "options": [
            "(1)   only (A) is true",
            "(2)   only (B) is true",
            "(3)   only (C) is true",
            "(4)   only (B) and (C) are true",
            "(5)   only (A) and (C) are true",
        ],
        "answer": "Ans: (4)",
    },
    {
        "number": "05",
        "question": (
            "A rocket, initially at rest, is fired vertically with an upward "
            "acceleration of 10 m s\u207b\u00b2. At an altitude of 0.5 km, the engine "
            "of the rocket cuts off. Its maximum altitude will be "
            "(Ignore the change in gravitational acceleration with altitude)"
        ),
        "options": [
            "(1)   1.9 km",
            "(2)   0.5 km",
            "(3)   1.5 km",
            "(4)   1.0 km",
            "(5)   1.6 km",
        ],
        "answer": "Ans: (4)",
    },
    {
        "number": "06",
        "question": (
            "A ball is projected horizontally with a velocity of 5 m s\u207b\u00b9 "
            "from the top of a building, 20 m high. The time the ball will take "
            "to hit the ground is"
        ),
        "options": [
            "(1)   \u221a2 s",
            "(2)   2 s",
            "(3)   \u221a3 s",
            "(4)   3 s",
            "(5)   5 s",
        ],
        "answer": "Ans: (2)",
    },
    {
        "number": "07",
        "question": (
            "A body displacement\u2013time graph is shown in the figure. "
            "(Refer to question paper for figure.)\n"
            "The corresponding speed (v)\u2013time (t) graph is best represented by:"
        ),
        "options": [
            "(1)   Graph 1",
            "(2)   Graph 2",
            "(3)   Graph 3",
            "(4)   Graph 4",
            "(5)   Graph 5",
        ],
        "answer": "Ans: (2)",
    },
    {
        "number": "08",
        "question": (
            "Following are displacement (s) vs time (t) graphs. "
            "(Refer to question paper for graphs A, B, C.)\n"
            "The graphs that show a positive acceleration, negative acceleration "
            "and zero acceleration respectively are,"
        ),
        "options": [
            "(1)   (A), (B), (C)",
            "(2)   (B), (C), (A)",
            "(3)   (C), (A), (D)",
            "(4)   (C), (B), (A)",
            "(5)   (A), (C), (B)",
        ],
        "answer": "Ans: (2)",
    },
    {
        "number": "09",
        "question": (
            "Consider the statements given about linear motion of an object:\n"
            "(A)  The displacement, velocity and acceleration always act in the same direction\n"
            "(B)  When the velocity is zero, displacement is also zero\n"
            "(C)  A velocity can exist where the acceleration is zero\n"
            "(D)  When an object moves backward, the displacement is always negative"
        ),
        "options": [
            "(1)   Only (A) is correct",
            "(2)   Only (B) is correct",
            "(3)   Only (C) is correct",
            "(4)   Only (C) and (D) are correct",
            "(5)   All (A), (B) and (C) are correct",
        ],
        "answer": "Ans: (3)",
    },
    {
        "number": "10",
        "question": (
            "Velocity time graph of a motion is shown in the figure. "
            "(Refer to question paper for figure.)\n"
            "Its displacement time graph is best represented by:"
        ),
        "options": [
            "(1)   Graph 1",
            "(2)   Graph 2",
            "(3)   Graph 3",
            "(4)   Graph 4",
            "(5)   Graph 5",
        ],
        "answer": "Ans: (5)",
    },
]

blank_layout = prs.slide_layouts[6]

for q in questions_data:
    slide = prs.slides.add_slide(blank_layout)

    # White background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE

    # ── Question number: top center, Red, Bell MT, 28pt ──────────────────────
    num_box = slide.shapes.add_textbox(
        Inches(0), Inches(0.15), Inches(13.33), Inches(0.7)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"Question  {q['number']}"
    run.font.name = "Bell MT"
    run.font.size = Pt(28)
    run.font.color.rgb = RED
    run.font.bold = True

    # ── Question text + options: Blue, Times New Roman, 36pt ─────────────────
    q_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.95), Inches(12.5), Inches(5.7)
    )
    tf2 = q_box.text_frame
    tf2.word_wrap = True

    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = q["question"]
    run2.font.name = "Times New Roman"
    run2.font.size = Pt(36)
    run2.font.color.rgb = BLUE

    # blank separator line
    tf2.add_paragraph()

    for opt in q["options"]:
        p_opt = tf2.add_paragraph()
        r = p_opt.add_run()
        r.text = opt
        r.font.name = "Times New Roman"
        r.font.size = Pt(36)
        r.font.color.rgb = BLUE

    # ── Answer: small text box, bottom-right, Black, Times New Roman, 36pt ───
    ans_box = slide.shapes.add_textbox(
        Inches(10.3), Inches(6.5), Inches(2.8), Inches(0.8)
    )
    tf3 = ans_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.RIGHT
    run3 = p3.add_run()
    run3.text = q["answer"]
    run3.font.name = "Times New Roman"
    run3.font.size = Pt(36)
    run3.font.color.rgb = BLACK
    run3.font.bold = True

prs.save("Kinematics_Model_Paper_2_QA.pptx")
print("Done! Saved as Kinematics_Model_Paper_2_QA.pptx")
